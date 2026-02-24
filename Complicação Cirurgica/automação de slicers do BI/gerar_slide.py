from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
import re

from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

EMU_PER_INCH = 914400
TITLE_P1 = "SINAIS DE COMPLICAÇÃO | RESULTADO POR ESTADO"
TITLE_P3 = "ORIENTAÇÕES DO CUIDADO PÓS CIRURGICO | RESULTADO POR ESTADO"


@dataclass
class SheetTable:
    name: str
    display_name: str
    headers: list[str]
    rows: list[list[object]]


def _inches_to_emu(value_in_inches: float) -> int:
    return int(Inches(value_in_inches))


def _emu_to_inches(value_emu: int) -> float:
    return float(value_emu) / EMU_PER_INCH


def _as_text(value: object, header: str) -> str:
    if value is None:
        return ""

    if isinstance(value, float) and value.is_integer():
        value = int(value)

    header_norm = header.upper()
    is_percent_col = (
        "PERCENTUAL" in header_norm
        or "PROPORCIONAL" in header_norm
        or "REPRESENTAT" in header_norm
    )

    if isinstance(value, (int, float)) and is_percent_col:
        num = float(value)
        if abs(num) <= 1:
            num *= 100
        return f"{num:.1f}%"

    if isinstance(value, int):
        return f"{value:,}".replace(",", ".")

    if isinstance(value, float):
        return f"{value:.1f}".replace(".", ",")

    return str(value)


def _display_name_from_sheet_name(sheet_name: str) -> str:
    cleaned = re.sub(r"^P\d+_", "", sheet_name, flags=re.IGNORECASE)
    cleaned = cleaned.replace("_", " ").strip()
    return cleaned if cleaned else sheet_name


def _find_primary_general_sheet(sheet_names: list[str]) -> str:
    if "GERAL" in sheet_names:
        return "GERAL"

    preferred = [name for name in sheet_names if name.upper() == "P1_GERAL"]
    if preferred:
        return preferred[0]

    general_candidates = [name for name in sheet_names if name.upper().endswith("_GERAL")]
    if general_candidates:
        return general_candidates[0]

    raise ValueError(
        'Nenhuma aba geral encontrada. Esperado "GERAL", "P1_GERAL" ou outro sufixo "_GERAL".'
    )


def _is_general_sheet(name: str) -> bool:
    return name.upper() == "GERAL" or name.upper().endswith("_GERAL")


def _is_p3_sheet(name: str) -> bool:
    return name.upper().startswith("P3_")


def _title_for_sheet(name: str) -> str:
    return TITLE_P3 if _is_p3_sheet(name) else TITLE_P1


def _extract_sheet_table(ws) -> SheetTable:
    raw_rows = list(ws.iter_rows(values_only=True))
    last_row_idx = 0
    last_col_idx = 0

    for r_idx, row in enumerate(raw_rows, start=1):
        row_has_value = False
        for c_idx, cell in enumerate(row, start=1):
            if cell is not None and str(cell).strip() != "":
                row_has_value = True
                if c_idx > last_col_idx:
                    last_col_idx = c_idx
        if row_has_value:
            last_row_idx = r_idx

    if last_row_idx == 0 or last_col_idx == 0:
        return SheetTable(
            name=ws.title,
            display_name=_display_name_from_sheet_name(ws.title),
            headers=["DADO"],
            rows=[],
        )

    matrix = []
    for row in raw_rows[:last_row_idx]:
        matrix.append(list(row[:last_col_idx]))

    headers = []
    for idx, h in enumerate(matrix[0], start=1):
        text = "" if h is None else str(h).strip()
        headers.append(text if text else f"COL_{idx}")

    data_rows = []
    for row in matrix[1:]:
        if any(cell is not None and str(cell).strip() != "" for cell in row):
            data_rows.append(row)

    return SheetTable(
        name=ws.title,
        display_name=_display_name_from_sheet_name(ws.title),
        headers=headers,
        rows=data_rows,
    )


def _set_slide_title(slide, text: str, slide_w: int, title_h: int, assets_dir: Path) -> None:
    side_margin = int(slide_w * 0.03)

    # Titulo
    text_left = side_margin + int(slide_w * 0.03)
    text_top = int(title_h * 0.08)
    text_w = int(slide_w * 0.75)
    text_h = int(title_h * 0.84)
    box = slide.shapes.add_textbox(text_left, text_top, text_w, text_h)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.bold = False
    run.font.name = "Calibri Light"
    run.font.color.rgb = RGBColor(32, 56, 100)  # #203864
    run.font.size = Pt(24)

    # Imagem no inicio, antes do titulo, centralizada verticalmente com a area do titulo
    inicio_path = assets_dir / "inicio.jpg"
    inicio_w = int(slide_w * 0.012)
    inicio_h = int(title_h * 0.60)
    inicio_x = side_margin
    inicio_y = int(text_top + (text_h - inicio_h) / 2)
    if inicio_path.exists():
        slide.shapes.add_picture(str(inicio_path), inicio_x, inicio_y, width=inicio_w, height=inicio_h)

    # Logo no canto superior direito
    logo_path = assets_dir / "logo.jpg"
    logo_w = int(slide_w * 0.045)
    logo_x = int(slide_w - side_margin - logo_w)
    logo_y = int(title_h * 0.12)
    if logo_path.exists():
        slide.shapes.add_picture(str(logo_path), logo_x, logo_y, width=logo_w)


def _column_weights(headers: list[str], rows: list[list[object]]) -> list[float]:
    widths = []
    for c, h in enumerate(headers):
        max_len = len(str(h))
        for row in rows:
            if c < len(row):
                max_len = max(max_len, len(_as_text(row[c], h)))
        widths.append(max(1.0, float(max_len)))
    return widths


def _fit_font_size_pt(table_h: int, row_count: int, slide_h: int) -> float:
    row_h = table_h / max(1, row_count)
    est_pt = (row_h / 12700.0) * 0.36
    slide_points = _emu_to_inches(slide_h) * 72
    min_pt = slide_points * 0.008
    max_pt = slide_points * 0.019
    return max(min_pt, min(est_pt, max_pt))


def _add_table_block(
    slide,
    table_data: SheetTable,
    block_left: int,
    block_top: int,
    block_w: int,
    block_h: int,
    slide_h: int,
    table_title_h: int | None = None,
) -> None:
    title_h = table_title_h if table_title_h is not None else int(block_h * 0.16)
    title_bottom_gap = int(block_h * 0.02)

    title_box = slide.shapes.add_textbox(block_left, block_top, block_w, title_h)
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = table_data.display_name
    run.font.bold = True
    run.font.color.rgb = RGBColor(31, 56, 100)
    slide_points = _emu_to_inches(slide_h) * 72
    run.font.size = Pt(slide_points * 0.022)

    table_top = block_top + title_h + title_bottom_gap
    raw_table_h = max(1, block_h - title_h - title_bottom_gap)
    # Compacta a area da tabela para reduzir altura visual das celulas.
    table_h = int(raw_table_h * 0.86)
    table_top += int((raw_table_h - table_h) / 2)

    row_count = len(table_data.rows) + 1
    col_count = max(1, len(table_data.headers))
    table_shape = slide.shapes.add_table(row_count, col_count, block_left, table_top, block_w, table_h)
    table = table_shape.table

    weights = _column_weights(table_data.headers, table_data.rows)
    weight_sum = sum(weights)
    for c in range(col_count):
        table.columns[c].width = int(block_w * (weights[c] / weight_sum))

    for r in range(row_count):
        table.rows[r].height = int(table_h / row_count)

    font_pt = _fit_font_size_pt(table_h, row_count, slide_h)

    for c, header in enumerate(table_data.headers):
        cell = table.cell(0, c)
        cell.text = str(header)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.runs[0]
        run.font.bold = True
        run.font.size = Pt(font_pt)
        run.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(38, 82, 181)

    for r, row in enumerate(table_data.rows, start=1):
        primeiro_valor = row[0] if len(row) > 0 else ""
        eh_total = str(primeiro_valor).strip().upper() == "TOTAL"
        for c in range(col_count):
            value = row[c] if c < len(row) else ""
            text = _as_text(value, table_data.headers[c])
            cell = table.cell(r, c)
            cell.text = text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(font_pt)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0, 0, 0)

            if eh_total:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(38, 82, 181)
                run.font.color.rgb = RGBColor(255, 255, 255)
            elif c == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(38, 82, 181)
                run.font.color.rgb = RGBColor(255, 255, 255)
            else:
                cell.fill.solid()
                if r % 2 == 0:
                    cell.fill.fore_color.rgb = RGBColor(237, 237, 237)
                else:
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)


def gerar_ppt(
    arquivo_excel: str = "indicadores_calculados.xlsx",
    arquivo_saida: str = "indicadores_apresentacao.pptx",
) -> None:
    assets_dir = Path(__file__).resolve().parent
    wb = load_workbook(arquivo_excel, data_only=True)
    ordered_names = wb.sheetnames
    primary_general = _find_primary_general_sheet(ordered_names)

    tables_by_sheet = {name: _extract_sheet_table(wb[name]) for name in ordered_names}

    prs = Presentation()
    prs.slide_width = _inches_to_emu(13.333)
    prs.slide_height = _inches_to_emu(7.5)

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 1a e 2a paginas em branco
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.slides.add_slide(prs.slide_layouts[6])

    # 3a pagina: GERAL
    geral_slide = prs.slides.add_slide(prs.slide_layouts[6])
    top_title_h = int(slide_h * 0.20)
    _set_slide_title(geral_slide, _title_for_sheet(primary_general), slide_w, top_title_h, assets_dir)

    right_half_left = int(slide_w * 0.50)
    right_half_w = int(slide_w * 0.50)
    right_margin = int(slide_w * 0.04)
    right_internal_margin = int(slide_w * 0.02)

    table_title_band_h = int(slide_h * 0.10)
    bottom_visible_gap = int(slide_h * 0.08)

    block_left = right_half_left + right_internal_margin
    block_w = right_half_w - right_margin - right_internal_margin
    block_top = top_title_h
    block_h = slide_h - top_title_h - bottom_visible_gap

    _add_table_block(
        geral_slide,
        tables_by_sheet[primary_general],
        block_left,
        block_top,
        block_w,
        block_h,
        slide_h,
        table_title_h=table_title_band_h,
    )

    # Demais abas apos a aba geral principal, na ordem original do Excel
    remaining_names = [name for name in ordered_names if name != primary_general]

    page_title_h = int(slide_h * 0.20)
    panel_top = page_title_h

    # Layout individual para abas gerais (ex.: P3_GERAL), igual ao slide da aba geral principal
    single_right_half_left = int(slide_w * 0.50)
    single_right_half_w = int(slide_w * 0.50)
    single_right_margin = int(slide_w * 0.04)
    single_right_internal_margin = int(slide_w * 0.02)
    single_table_title_band_h = int(slide_h * 0.10)
    single_bottom_visible_gap = int(slide_h * 0.08)
    single_block_left = single_right_half_left + single_right_internal_margin
    single_block_w = single_right_half_w - single_right_margin - single_right_internal_margin
    single_block_top = page_title_h
    single_block_h = slide_h - page_title_h - single_bottom_visible_gap

    # Layout de detalhamento: 4 por slide (2 colunas x 2 linhas)
    grid_cols = 2
    grid_rows = 2
    grid_side_margin = int(slide_w * 0.04)
    grid_col_gap = int(slide_w * 0.03)
    grid_bottom_margin = int(slide_h * 0.12)
    grid_row_gap = int(slide_h * 0.04)
    grid_total_h = slide_h - page_title_h - grid_bottom_margin
    grid_panel_w = int((slide_w - (2 * grid_side_margin) - ((grid_cols - 1) * grid_col_gap)) / grid_cols)
    grid_panel_h = int((grid_total_h - ((grid_rows - 1) * grid_row_gap)) / grid_rows)

    i = 0
    while i < len(remaining_names):
        current_name = remaining_names[i]

        # Toda aba *_GERAL deve ficar sozinha em uma pagina individual
        if _is_general_sheet(current_name):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            _set_slide_title(slide, _title_for_sheet(current_name), slide_w, page_title_h, assets_dir)
            _add_table_block(
                slide,
                tables_by_sheet[current_name],
                single_block_left,
                single_block_top,
                single_block_w,
                single_block_h,
                slide_h,
                table_title_h=single_table_title_band_h,
            )
            i += 1
            continue

        # Coleta ate 4 abas nao-gerais preservando ordem e sem cruzar uma aba geral.
        chunk: list[str] = []
        while i < len(remaining_names) and len(chunk) < 4 and not _is_general_sheet(remaining_names[i]):
            chunk.append(remaining_names[i])
            i += 1

        if not chunk:
            continue

        # Se houver qualquer aba P3 no bloco, usa titulo de P3.
        slide_title = TITLE_P3 if any(_is_p3_sheet(name) for name in chunk) else TITLE_P1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        _set_slide_title(slide, slide_title, slide_w, page_title_h, assets_dir)

        for idx, sheet_name in enumerate(chunk):
            row = idx // grid_cols
            col = idx % grid_cols
            x = grid_side_margin + col * (grid_panel_w + grid_col_gap)
            y = panel_top + row * (grid_panel_h + grid_row_gap)
            _add_table_block(
                slide,
                tables_by_sheet[sheet_name],
                x,
                y,
                grid_panel_w,
                grid_panel_h,
                slide_h,
            )

    prs.save(arquivo_saida)


if __name__ == "__main__":
    gerar_ppt(
        arquivo_excel="indicadores_calculados_formatado.xlsx",
        arquivo_saida="indicadores_apresentacao.pptx",
    )
    print("Arquivo gerado: indicadores_apresentacao.pptx")
